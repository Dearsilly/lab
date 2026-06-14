"""Generate defense PPT for AES system — Web + Android."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
C_DARK = RGBColor(0x1F, 0x38, 0x64)
C_BLUE = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT = RGBColor(0xD5, 0xE8, 0xF0)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x6B, 0x72, 0x80)
C_GREEN = RGBColor(0x10, 0xB9, 0x81)
C_INDIGO = RGBColor(0x63, 0x66, 0xF1)
C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
C_DARKTEXT = RGBColor(0x1F, 0x29, 0x37)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text, subtitle=""):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

def add_footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "AES 中英文双语作文评分系统  |  工程实践项目答辩  |  苏州 · 2026"
    p.font.size = Pt(10)
    p.font.color.rgb = C_GRAY

def add_bullets(slide, left, top, width, height, items, size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = C_DARKTEXT
        p.space_after = Pt(8)

def add_card(slide, left, top, width, height, title, content, color=C_INDIGO):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    for line in content:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(14)
        p2.font.color.rgb = C_DARKTEXT
        p2.space_after = Pt(4)

def add_layer_card(slide, left, top, width, title, desc, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_GRAY

# ====== Slide 1: Title ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "中英文双语作文自动评分系统"
p.font.size = Pt(44)
p.font.color.rgb = C_WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Automated Essay Scoring with BERT · Web + Android"
p2.font.size = Pt(22)
p2.font.color.rgb = C_LIGHT
p2.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10), Inches(1.5))
tf2 = txBox2.text_frame
for line in ["4 人团队  |  工程实践项目  |  苏州  |  2026"]:
    p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# ====== Slide 2: Problem ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "问题背景", "为什么需要自动作文评分？")
add_footer(slide)
add_bullets(slide, 0.8, 1.6, 11, 5.5, [
    "▸ 教师批改痛点：4 个班 × 50 篇 = 200 篇作文，每篇 10 分钟 = 33 小时工作量",
    "▸ 评分一致性：疲劳导致前 30 篇精批细改，后面越批越潦草",
    "▸ 缺乏分项反馈：学生只知道总分，不知道内容/结构/语言哪个弱",
    "▸ 商业方案昂贵：ETS e-rater 等闭源，不适用于中文教学场景",
    "",
    "→ 目标：开源、中英双语、Web + Android 双平台的自动评分系统，",
    "   让教师和学生都能随时随地获取即时评分和反馈"
])

# ====== Slide 3: System Overview ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "系统功能概览")
add_footer(slide)
add_card(slide, 0.5, 1.5, 3.8, 2.5, "单篇评分", [
    "中英文自动检测",
    "语言手动切换",
    "总分仪表盘 + 维度雷达图",
    "模板化评语反馈"
], C_INDIGO)
add_card(slide, 4.6, 1.5, 3.8, 2.5, "批量评分", [
    "CSV 文件导入",
    "逐篇推理 + 进度展示",
    "结果列表 + 下载导出",
    "最多 100 篇"
], C_BLUE)
add_card(slide, 8.7, 1.5, 3.8, 2.5, "中英对比", [
    "同一文本双模型评分",
    "中英模型并排展示",
    "维度分差直观对比",
    "Web + Android 双端"
], C_GREEN)
add_card(slide, 0.5, 4.3, 3.8, 2.8, "双平台", [
    "Web: Flask + Streamlit",
    "Android: Kotlin + Compose",
    "Android 完全本地推理",
    "PyTorch Mobile · 离线可用"
], C_ORANGE)
add_card(slide, 4.6, 4.3, 7.9, 2.8, "核心技术栈", [
    "Python / Kotlin    |    BERT 双模型    |    MSE + AdamW + fp16",
    "Flask REST API    |    Streamlit + Plotly    |    Jetpack Compose + Canvas",
    "PyTorch Mobile Lite    |    Kotlin WordPiece 分词器    |    pytest + Playwright"
], C_DARK)

# ====== Slide 4: Architecture ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "系统架构", "分层设计 · Web + Android 双端共享模型")
add_footer(slide)
add_bullets(slide, 0.8, 1.4, 5.5, 5.5, [
    "▸ 数据层",
    "    ASAP 英文 12,979 篇 + Google 翻译中文 3,950 篇",
    "    按 essay_set 级 Prompt 隔离划分训练/测试集",
    "",
    "▸ 模型层",
    "    BERT-base-uncased (EN) + BERT-base-chinese (ZH)",
    "    [CLS] → Dropout → Linear → Sigmoid → [0,1]",
    "    多任务变体：DeBERTa + 4 回归头（已实现）",
    "",
    "▸ 推理引擎层",
    "    语言检测 → 分词 → 模型路由 → 推理 → 反馈",
    "    Android: Kotlin 全自主实现（零 NLP 库依赖）",
    "",
    "▸ 应用层",
    "    Web: Flask API + Streamlit UI (4 页面)",
    "    Android: Compose + Bottom Nav (4 Tab)"
], size=16)
add_layer_card(slide, 7.3, 1.5, 5.2, "Web 版", "Flask API → GPU 推理 → ~200ms/篇", C_INDIGO)
add_layer_card(slide, 7.3, 2.6, 5.2, "Android 版", "PyTorch Mobile → CPU 推理 → ~1-5s/篇", C_GREEN)
add_layer_card(slide, 7.3, 3.7, 5.2, "共 享", "同一套 .pt 模型权重，双平台部署", C_ORANGE)
# Comparison table
add_bullets(slide, 7.3, 4.7, 5.2, 2.5, [
    "  Web: 分词 HuggingFace / 语言 langdetect",
    "  Android: 分词 Kotlin 自实现 / 语言字符规则",
    "  Web: GPU fp16 / Android: CPU fp32",
    "  Web: 需要网络 / Android: 完全离线",
], size=13)

# ====== Slide 5: BERT Model ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "BERT 回归模型", "[CLS] → Dropout → Linear → Sigmoid")
add_footer(slide)
add_bullets(slide, 0.8, 1.5, 5.5, 5, [
    "▸ 英文模型: bert-base-uncased",
    "   12 层 Transformer, 768 维, 110M 参数",
    "   WordPiece 词表 30,522 tokens",
    "   fp16 混合精度训练, batch_size=16",
    "",
    "▸ 中文模型: bert-base-chinese",
    "   同架构, 词表 21,128 tokens",
    "   jieba 预分词 + 全角转半角",
    "   翻译数据 3,950 篇, batch_size=8",
    "",
    "▸ 训练配置",
    "   AdamW lr=2e-5 | Warmup 10% | 早停 patience=3",
    "   MSE Loss | 5 epochs | English fp16, Chinese fp32"
], size=16)
add_card(slide, 7.5, 1.5, 5, 2.5, "模型表现", [
    "BERT-base 英文: Val QWK 0.58",
    "BERT-base 中文: Val QWK 0.79, Test 0.76",
    "RoBERTa-base: Val QWK 0.58 (备选)",
    "",
    "英文 QWK 低 = 严格的 Prompt 隔离评估",
    "中文 QWK 高 = 翻译同源 + 随机划分",
], C_INDIGO)
add_card(slide, 7.5, 4.3, 5, 2.8, "多任务扩展 (已实现)", [
    "DeBERTa-v3 + 4 个独立回归头",
    "Head 0: 总评分    Head 2: 结构",
    "Head 1: 内容       Head 3: 语言",
    "训练策略: 总评分权重 1.0",
    "         维度权重各 0.5",
    "额外显存 < 1MB, 待部署",
], C_ORANGE)

# ====== Slide 6: Data Pipeline ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "数据处理流水线", "Prompt 隔离划分 — 防止数据泄露的关键设计")
add_footer(slide)
add_card(slide, 0.8, 1.5, 5.5, 5.2, "英文数据流", [
    "ASAP CSV (12,976 篇, 8 个题目集)",
    "↓ HTML unescape, URL/@mention 移除",
    "↓ 空白规范化 (\\s+ → 空格)",
    "↓ 按 essay_set 归一化 [0,1]",
    "↓ BERT 分词 (max_length=512)",
    "↓ essay_set 隔离划分 → train/val/test",
    "↓ 断言验证: train ∩ test = ∅",
], C_INDIGO)
add_card(slide, 7.5, 1.5, 5.5, 5.2, "中文数据流 (差异点)", [
    "ASAP CSV → Google Translate (en→zh)",
    "↓ 全角转半角, 控制字符移除",
    "↓ jieba 分词 (词间加空格)",
    "↓ 按 essay_set 归一化 [0,1]",
    "↓ BERT 分词 (bert-base-chinese)",
    "↓ 随机划分 (翻译同源, 隔离无意义)",
    "↓ 英文 QWK=0.58, 中文 QWK=0.79 ← 差异原因",
], C_GREEN)

# ====== Slide 7: Android Architecture ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "Android 版架构", "Kotlin + Jetpack Compose + PyTorch Mobile · 完全离线")
add_footer(slide)
add_card(slide, 0.5, 1.5, 3.8, 2.5, "UI 层", [
    "Jetpack Compose + Material3",
    "Bottom Navigation (4 Tab)",
    "Canvas 自绘雷达图",
    "评分 / 批量 / 对比 / 设置",
], C_INDIGO)
add_card(slide, 4.6, 1.5, 3.8, 2.5, "推理层", [
    "AESPredictor (模型管理)",
    "BertTokenizer (100行 Kotlin)",
    "LanguageDetector (字符规则)",
    "PyTorch Mobile Lite 1.13.1",
], C_GREEN)
add_card(slide, 8.7, 1.5, 3.8, 2.5, "assets/", [
    "bert_model.pt (418MB)",
    "zh_model.pt   (391MB)",
    "vocab.txt     (30,522)",
    "zh_vocab.txt  (21,128)",
], C_ORANGE)
add_card(slide, 0.5, 4.3, 12, 2.8, "关键技术决策", [
    "分词器: Kotlin 自实现 WordPiece (贪心最长子词匹配 + ##前缀), 零 Java NLP 依赖",
    "语言检测: 统计 Unicode 范围 (一-鿿) 中文字符占比 >50% → zh, 零外部库",
    "雷达图: Compose Canvas 自绘 (网格 + 轴线 + 数据多边形 + 标签), 无第三方图表库",
    "状态管理: remember + mutableStateOf, Coroutines(Default) 推理, withContext(Main) 更新 UI",
], C_DARK)

# ====== Slide 8: Android Screens ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "Android 界面展示", "4 个 Tab · 中文 UI")
add_footer(slide)
add_card(slide, 0.5, 1.5, 3.8, 3, "评分 Tab", [
    "语言 Chip: 自动/英文/中文",
    "文本输入 (≤10000 字符)",
    "实数字符/词数/语言统计",
    "总分仪表盘 + 渐变进度条",
    "三维度卡片 + 雷达图 + 评语"
], C_INDIGO)
add_card(slide, 4.6, 1.5, 3.8, 3, "批量 Tab", [
    "CSV 文件选择器",
    "逐条推理 + 进度条",
    "LazyColumn 结果列表",
    "显示 id/分数/状态"
], C_BLUE)
add_card(slide, 8.7, 1.5, 3.8, 3, "对比 + 设置", [
    "中英双模型 → 左右并排",
    "总分 + 三维度差值对比",
    "模型状态卡片 (✅/❌)",
    "QWK 信息 · 重新加载按钮"
], C_GREEN)
add_bullets(slide, 0.5, 4.8, 12, 2.5, [
    "所有组件自主实现: RadarChart (Canvas) / ScoreGauge (渐变进度条) / FeedbackCard (彩色边框)",
    "总计约 800 行 Kotlin · 最低 SDK 29 (Android 10) · Compose Navigation 单 Activity 架构"
], size=16)

# ====== Slide 9: Kotlin Tokenizer ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "Kotlin BERT 分词器", "100 行代码 · 零外部 NLP 库 · 自主实现")
add_footer(slide)
add_bullets(slide, 0.8, 1.5, 6, 5, [
    "▸ 词表加载: 从 assets/vocab.txt 读取 → HashMap<String, Int>",
    "    英文: 30,522 tokens (bert-base-uncased)",
    "    中文: 21,128 tokens (bert-base-chinese)",
    "",
    "▸ basicTokenize: 基础分词",
    "    中文: 逐字切分 (Unicode 范围 一-鿿)",
    "    英文: 按空格/标点拆分 + 小写化",
    "",
    "▸ wordPieceTokenize: 子词匹配",
    "    贪心最长匹配: 从词尾向前搜索",
    "    首次不加前缀, 后续加 ## 前缀",
    "    未匹配 → 使用 [UNK] 回退",
    "",
    "▸ tokenize: 完整流程",
    "    文本 → basicTokenize → wordPieceTokenize",
    "    → [CLS] + tokens + [SEP] → LongArray(512) 填充",
], size=16)
add_card(slide, 7.5, 1.5, 5, 5.5, "关键代码 (~20行核心)", [
    "fun wordPieceTokenize(word): List<Int> {",
    "  var remaining = word, isFirst = true",
    "  while (remaining.isNotEmpty()) {",
    "    prefix = if (isFirst) '' else '##'",
    "    // 从最长匹配递减搜索",
    "    for (end in remaining.length downTo 1) {",
    "      id = vocab[prefix+remaining[0:end]]",
    "      if (id != null) { tokens.add(id); break }",
    "    }",
    "    if (!found) { tokens.add([UNK]); break }",
    "    isFirst = false",
    "  }",
    "}",
    "",
    "为什么自实现？",
    "• HuggingFace Tokenizers 无 Android 版",
    "• Java NLP 库 (DJL) 体积大、依赖多",
    "• 100 行代码完全可控，便于调试修改",
], C_DARK)

# ====== Slide 10: Training & Evaluation ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "训练与评估", "QWK — 衡量人机评分一致性的主要指标")
add_footer(slide)
add_bullets(slide, 0.8, 1.5, 5.5, 3, [
    "▸ QWK (Quadratic Weighted Kappa)",
    "    二次权重矩阵: w_ij = (i-j)^2",
    "    评分偏差越大, 惩罚越重",
    "    范围 [-1, 1], ≥0.70 为可接受水平",
    "",
    "▸ 辅助指标: MAE (平均绝对误差)",
    "    Pearson r (线性相关系数)"
], size=16)

# Results table
add_bullets(slide, 0.8, 4.2, 5.5, 2.5, [
    " 模型             Val QWK   Test QWK",
    " BERT-base EN       0.58        —",
    " BERT-base ZH       0.79       0.76",
    " RoBERTa EN         0.58        —",
    "",
    " 英文 0.58 = 严格 Prompt 隔离下的真实泛化",
    " 中文 0.79 = 翻译同源 + 随机划分 (虚高)",
], size=15)
add_card(slide, 7.5, 1.5, 5, 5.5, "测试策略", [
    "Web 版单元测试 (pytest): 17 项",
    "  数据预处理 · API 端点 · 评估指标",
    "  文本清洗/归一化/QWK正确性",
    "",
    "Web 版 E2E 测试 (Playwright): 41 项",
    "  API 全端点: Health/Score/Batch/404",
    "  UI 全页面: 评分/批量/对比/设置",
    "  范文验证: 中英文高分范文评分",
    "  错误路径: 空提交/超长/格式错误",
    "",
    "Android 版: 手动测试 (本地演示)",
], C_GREEN)

# ====== Slide 11: Tech Decisions ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "关键技术决策", "10 个架构决策及其理由")
add_footer(slide)
decisions = [
    ("BERT-base 优先", "8GB 显存稳定, 生态成熟", C_INDIGO),
    ("Prompt 隔离划分", "防止题目级数据泄露, 评估更真实", C_BLUE),
    ("中文翻译数据", "快速概念验证, 绕过数据获取瓶颈", C_GREEN),
    ("Flask + Streamlit", "团队技能 Python 全栈, 快速迭代", C_INDIGO),
    ("PyTorch Mobile", "Android 本地推理, 离线可用", C_GREEN),
    ("Kotlin 分词器", "避免 Java NLP 库依赖, 100行可控", C_ORANGE),
    ("Canvas 雷达图", "无第三方图表库, Android 原生渲染", C_INDIGO),
    ("无 DB/Docker/认证", "控制工程实践复杂度边界", C_BLUE),
    ("MSE + Sigmoid", "训练稳定, 梯度有上界, 输出天然 [0,1]", C_GREEN),
    ("字符规则语言检测", "Android 端零依赖, 准确率足够 (>95%)", C_ORANGE),
]
for i, (title, reason, color) in enumerate(decisions):
    y = 1.5 + i * 0.55
    shape = slide.shapes.add_shape(1, Inches(1.5), Inches(y), Inches(10.3), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = f"▸ {title}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = reason
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_GRAY

# ====== Slide 12: Summary ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_title_bar(slide, "项目成果总结")
add_footer(slide)
add_card(slide, 0.5, 1.5, 6, 3.5, "已完成", [
    "✅ 中英文双语 BERT 评分 (双模型)",
    "✅ 多维度评分 + 雷达图 + 评语反馈",
    "✅ Web 版: Flask API + Streamlit UI",
    "✅ Android 版: Compose + PyTorch Mobile",
    "✅ Kotlin 自实现 BERT 分词器",
    "✅ Android 完全离线本地推理",
    "✅ 批量评分 CSV 流水线",
    "✅ 17 单元测试 + 41 E2E 测试",
], C_GREEN)
add_card(slide, 6.8, 1.5, 6, 3.5, "后续改进", [
    "🔮 真实中文作文数据训练",
    "🔮 多任务 DeBERTa-v3 部署",
    "🔮 Longformer 长文本支持",
    "🔮 LLM 个性化评语生成",
    "🔮 模型版本管理 + 持续评估",
    "🔮 Android APK 体积优化 (按需下载)",
    "🔮 用户认证 + 历史记录",
    "🔮 iOS 版 (Swift + CoreML)",
], C_ORANGE)
add_card(slide, 0.5, 5.3, 12, 1.8, "核心创新", [
    "中英文双模型 + 双平台 (Web+Android)  |  Kotlin 自主实现 BERT 分词器  |  Android 完全离线推理",
    "Compose Canvas 自绘雷达图  |  Prompt 隔离严格评估  |  58 项自动化测试覆盖"
], C_DARK)

# ====== Slide 13: Thank You ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢！欢迎提问"
p.font.size = Pt(48)
p.font.color.rgb = C_WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Thank You · Questions Welcome"
p2.font.size = Pt(24)
p2.font.color.rgb = C_LIGHT
p2.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10), Inches(1.5))
tf2 = txBox2.text_frame
for line in ["中英文双语作文自动评分系统 (AES)", "BERT + PyTorch Mobile · Web + Android", "工程实践项目 · 苏州 · 2026"]:
    p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER

# Save
out = "doc/AES_答辩PPT.pptx"
prs.save(out)
print(f"Done: {out} ({len(prs.slides)} slides)")
