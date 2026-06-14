#!/usr/bin/env python3
"""AES 自动作文评分系统 — 答辩 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = 'doc/AES_答辩PPT.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── Color Palette ─────────────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x56, 0xDB)  # 深蓝
SECONDARY  = RGBColor(0x37, 0x41, 0x51)  # 深灰
ACCENT     = RGBColor(0x10, 0xB9, 0x81)  # 绿色
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF8, 0xF9, 0xFA)
BLACK      = RGBColor(0x00, 0x00, 0x00)

# ── Helpers ───────────────────────────────────────────────────────

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_title_bar(slide, title_text):
    """Add a dark blue title bar at top"""
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), PRIMARY)
    tf = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9)).text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True

def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.6), width=Inches(11.5), height=Inches(5.2), size=Pt(22)):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = SECONDARY
    p.space_after = Pt(12)
    return tf

def add_bullet_list(slide, items, left=Inches(0.8), top=Inches(1.6), width=Inches(11.5), height=Inches(5.2), size=Pt(20)):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = SECONDARY
        p.space_after = Pt(10)
        p.level = 0
    return tf

def add_simple_table(slide, headers, rows, left=Inches(1.5), top=Inches(2.0), width=Inches(10.3), row_height=Inches(0.55)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows).table

    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = SECONDARY
                p.alignment = PP_ALIGN.CENTER
    return table

# ═══════════════════════════════════════════════════════════════════════
# Slide 1: Title Slide
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s, PRIMARY)
tf = s.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(1.5)).text_frame
p = tf.paragraphs[0]
p.text = 'AES 自动作文评分系统'
p.font.size = Pt(48)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = 'Automated Essay Scoring System'
p2.font.size = Pt(28)
p2.font.color.rgb = RGBColor(0xBB, 0xCA, 0xE8)
p2.alignment = PP_ALIGN.CENTER

tf2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(2.5)).text_frame
for line in ['基于深度学习的英文/中文作文智能评分', '', '工程实践项目  2026年6月']:
    p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════
# Slide 2: Agenda
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '汇报大纲')
items = [
    '1.  项目背景 —— 为什么要做自动作文评分？',
    '2.  系统架构 —— Web版 + Android版 双端架构',
    '3.  技术方案 —— BERT 模型设计、数据处理与训练策略',
    '4.  实现细节 —— 推理引擎、API、UI、Android分词器',
    '5.  测试结果 —— 模型表现、单元测试与E2E测试',
    '6.  总结与展望 —— 成果、限制与改进方向',
]
add_bullet_list(s, items, top=Inches(1.8), height=Inches(5.0), size=Pt(24))

# ═══════════════════════════════════════════════════════════════════════
# Slide 3: Problem & Motivation
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '1. 项目背景与问题定义')
items = [
    '• 作文评分是教师最耗时的工作：4个班 × 50篇 × 5分钟 = 16+小时',
    '• 传统自动评分依赖手工特征（词汇、句长、语法），难以捕捉深层语义',
    '• 现有深度学习方案大多仅支持英文，缺乏全平台覆盖',
    '',
    '项目目标：',
    '• 构建中英文双语BERT作文自动评分系统',
    '• 提供多维度评分（总分 + 内容/结构/语言）和文字评语反馈',
    '• 实现 Web版 和 Android本地推理版 双版本完整前端',
]
add_bullet_list(s, items, top=Inches(1.6), size=Pt(22))

# ═══════════════════════════════════════════════════════════════════════
# Slide 4: System Architecture
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '2. 系统架构 — 三层分层设计')
items = [
    '数据层：ASAP 英文作文数据集（12,979篇）+ Google翻译中文版（3,950篇）',
    '  ├─ 英文：HTML清理 → 归一化 → prompt隔离划分',
    '  └─ 中文：全半角转换 → jieba分词 → 随机划分',
    '',
    '模型层：BERT-base-uncased（英文）+ BERT-base-chinese（中文）',
    '  ├─ [CLS] → Dropout(0.1) → Linear(768→1) → Sigmoid → [0,1]',
    '  └─ 预留多任务变体（4回归头：总分/内容/结构/语言）',
    '',
    '应用层：三端统一API设计，4个REST端点',
    '  ├─ Flask REST API（后端服务）',
    '  ├─ Streamlit Web UI（4页面：评分/批量/对比/设置）',
    '  └─ Android App（PyTorch Mobile本地推理，4 Tab导航）',
]
add_bullet_list(s, items, top=Inches(1.6), size=Pt(19))

# ═══════════════════════════════════════════════════════════════════════
# Slide 5: Model Architecture
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '3. 模型架构 — BERT + 回归头')
# Architecture diagram using text + shapes
tf = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2)).text_frame
tf.word_wrap = True

steps = [
    ('输入文本', 'English text / 中文文本'),
    ('BERT Tokenizer', 'WordPiece (30522 EN / 21128 ZH vocab)'),
    ('BERT Encoder', '12层 Transformer, 768-d hidden, 110M params'),
    ('[CLS] Pooling', '取第一个 token 的输出向量 (768-d)'),
    ('Dropout (0.1)', ''),
    ('Linear (768 → 1)', ''),
    ('Sigmoid', '输出分数 ∈ [0, 1]'),
]

for i, (title, desc) in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f'{title}'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(2)
    if desc:
        p2 = tf.add_paragraph()
        p2.text = f'      └─ {desc}'
        p2.font.size = Pt(14)
        p2.font.color.rgb = SECONDARY
        p2.space_after = Pt(4)
    # arrow
    if i < len(steps) - 1:
        pa = tf.add_paragraph()
        pa.text = '      ↓'
        pa.font.size = Pt(12)
        pa.font.color.rgb = ACCENT
        pa.space_after = Pt(4)

# ═══════════════════════════════════════════════════════════════════════
# Slide 6: Data Pipeline
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '3. 数据处理与训练策略')

add_simple_table(s,
    ['参数', '英文模型', '中文模型'],
    [
        ['基础模型', 'bert-base-uncased', 'bert-base-chinese'],
        ['训练数据', 'ASAP 12,979篇', '翻译数据 3,950篇'],
        ['学习率', '2e-5', '2e-5'],
        ['Batch Size', '16', '8'],
        ['Epochs', '5（早停 patience=3）', '5（早停 patience=3）'],
        ['精度', 'fp16 混合精度', 'fp32'],
        ['优化器', 'AdamW (wd=0.01)', 'AdamW (wd=0.01)'],
        ['划分策略', 'Essay-set 级别隔离', '随机划分'],
    ],
    top=Inches(1.6)
)

# ═══════════════════════════════════════════════════════════════════════
# Slide 7: Web vs Android
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '4. 实现细节 — Web版与Android版对比')

add_simple_table(s,
    ['维度', 'Web 版', 'Android 版'],
    [
        ['推理位置', 'Flask 服务器 (GPU fp16)', '手机本地 (CPU)'],
        ['模型加载', 'Python PyTorch', 'PyTorch Mobile Lite 1.13.1'],
        ['分词器', 'HuggingFace Transformers', 'Kotlin 自实现 WordPiece'],
        ['语言检测', 'langdetect + 字符规则', '字符规则（零依赖）'],
        ['UI 框架', 'Streamlit + Plotly', 'Jetpack Compose + Canvas'],
        ['网络', '需要 HTTP 连接', '完全离线'],
        ['推理速度', '~200ms', '~1-5s'],
    ],
    top=Inches(1.6)
)

# ═══════════════════════════════════════════════════════════════════════
# Slide 8: Android Architecture
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '4. Android App 架构')

items = [
    'UI层 (Jetpack Compose + Material3)',
    '  ┌──────┐  ┌──────┐  ┌──────┐  ┌──────┐',
    '  │ 评分  │  │ 批量  │  │ 对比  │  │ 设置  │  底部4 Tab NavigationBar',
    '  └──────┘  └──────┘  └──────┘  └──────┘',
    '',
    '推理层 (inference/)',
    '  ├─ AESPredictor: 模型加载、推理调度、反馈生成',
    '  ├─ BertTokenizer: Kotlin原生WordPiece分词器（~100行代码）',
    '  └─ LanguageDetector: 中文字符统计规则（Unicode 一-鿿）',
    '',
    '本地资源 (assets/)',
    '  ├─ bert_model.pt（英文，418MB）+ zh_model.pt（中文，391MB）',
    '  └─ vocab.txt（英文词表，30522 tokens）+ zh_vocab.txt（中文词表，21128 tokens）',
    '',
    '核心技术栈：Kotlin 2.2 | Compose BOM 2024.09 | PyTorch Mobile 1.13.1 | minSdk 29',
]
add_bullet_list(s, items, top=Inches(1.5), size=Pt(17))

# ═══════════════════════════════════════════════════════════════════════
# Slide 9: Tokenizer (Android highlight)
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '4. Kotlin BERT WordPiece 分词器实现')

items = [
    'Android版核心技术亮点：Kotlin 原生 BERT WordPiece 分词器',
    '',
    '分词流程：',
    '  1. basicTokenize(): 按空格/标点分词，中文字符逐字切分',
    '  2. wordPieceTokenize(): 贪心最长子词匹配，##前缀标记子词',
    '  3. 序列组装: [CLS] + tokens + [SEP] → LongArray(512)',
    '',
    '实现特点：',
    '  • 约100行Kotlin代码，零外部Java NLP库依赖',
    '  • 词表从 assets/vocab.txt 读取（每行一个token），构建 HashMap<String, Int>',
    '  • 中文检测通过Unicode范围 一-鿿 统计中文字符占比（>50% → 中文）',
    '  • 正确的截断保护：remaining = maxLength - 1 - tokens.size，保留[SEP]空间',
    '',
    '性能：在手机CPU上，分词+推理总耗时约1-5秒（取决于设备性能）',
]
add_bullet_list(s, items, top=Inches(1.5), size=Pt(17))

# ═══════════════════════════════════════════════════════════════════════
# Slide 10: Model Performance
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '5. 模型评估结果')

add_simple_table(s,
    ['模型', 'Val QWK', 'Test QWK', '评估方式', '备注'],
    [
        ['BERT-base 英文', '0.58', '—', 'prompt 隔离', '更真实的跨题目泛化'],
        ['BERT-base 中文', '0.79', '0.76', '随机划分', '翻译数据同源'],
    ],
    top=Inches(1.8)
)

tf = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(2.5)).text_frame
tf.word_wrap = True
for i, line in enumerate([
    '• 英文 QWK 0.58 = prompt隔离条件下的真实泛化能力（训练/测试包含不同题目）',
    '• 中文 QWK 0.79 = 随机划分，偏高但证明模型可从翻译数据中学习',
    '• QWK ≥ 0.70 为可接受水平，英文模型在严格评估下仍有提升空间',
    '• 辅助指标：MAE（平均绝对误差）、Pearson 相关系数',
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = SECONDARY
    p.space_after = Pt(6)

# ═══════════════════════════════════════════════════════════════════════
# Slide 11: Testing
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '5. 测试策略')

items = [
    '单元测试（17项 — pytest）',
    '  • 数据预处理: cleaner, normalizer, splitter, tokenizer',
    '  • API端点: /health, /models, /score, /batch',
    '  • 覆盖正常输入和边界条件（空文本、超长文本、缺失参数）',
    '',
    '端到端测试（41项 — Playwright）',
    '  • 自动启动 Flask + Streamlit 服务',
    '  • 在无界面浏览器中执行完整用户操作流程',
    '  • 覆盖 Flask API 全部 4 端点 + Streamlit UI 全部 4 页面',
    '',
    '运行：pytest tests/ -v （单元）| python tests/test_e2e.py （E2E）',
]
add_bullet_list(s, items, top=Inches(1.6), size=Pt(20))

# ═══════════════════════════════════════════════════════════════════════
# Slide 12: Technical Decisions
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '技术决策记录')

add_simple_table(s,
    ['决策项', '选择', '理由'],
    [
        ['主模型', 'BERT-base (110M)', '8GB显存稳定运行，生态成熟'],
        ['中文数据', 'Google Translate 翻译', '快速启动，概念验证'],
        ['分数归一化', '按 essay_set 归一化', '保持各题目间的可比性'],
        ['数据划分', 'Prompt 隔离（英文）', '防止题目级数据泄露'],
        ['Web技术栈', 'Flask + Streamlit', '团队技能栈匹配，快速开发'],
        ['移动端推理', 'PyTorch Mobile 本地', '离线可用，无需网络和服务器'],
        ['复杂度边界', '不引入数据库/认证', '聚焦核心评分流水线'],
    ],
    top=Inches(1.6)
)

# ═══════════════════════════════════════════════════════════════════════
# Slide 13: Summary
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, '6. 项目成果总结')

items = [
    '1.  中英文双语BERT自动作文评分模型：英文QWK 0.58 / 中文QWK 0.79',
    '2.  完整的Flask REST API + Streamlit Web UI（4页面，中英文切换）',
    '3.  Android原生App：PyTorch Mobile本地推理，离线可用',
    '4.  Kotlin BERT WordPiece分词器：从零实现，零外部NLP依赖',
    '5.  17项单元测试 + 41项端到端测试覆盖',
    '6.  完整技术文档（架构设计、答辩脚本、Word报告、PPT）',
    '',
    '已知限制：',
    '• 中文模型使用翻译数据 → 计划收集真实中文作文',
    '• 维度评分基于启发式拆分 → 预留多任务模型架构',
    '• BERT 512 token限制 → 考虑引入Longformer',
    '• Android双模型~800MB → 考虑INT8量化',
]
add_bullet_list(s, items, top=Inches(1.5), size=Pt(20))

# ── Save ────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f'Done: {OUTPUT_PATH}')
